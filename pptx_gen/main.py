import os
import shutil

from PIL import Image
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches


def combine_images_horizontal(images):
    """
    Объединяет несколько изображений типа PpmImageFile в одно горизонтальное изображение.

    Args:
      images: Список изображений типа PpmImageFile.

    Returns:
      Объединенное изображение типа PpmImageFile.
    """

    # Получаем размеры изображений
    widths = [image.width for image in images]
    heights = [image.height for image in images]

    # Вычисляем ширину и высоту конечного изображения
    total_width = sum(widths)
    max_height = max(heights)

    # Создаем новое изображение с нужными размерами
    new_image = Image.new("RGB", (total_width, max_height))

    # Вставляем изображения в новое изображение
    x_offset = 0
    for image in images:
        new_image.paste(image, (x_offset, 0))
        x_offset += image.width

    return new_image


def create_presentation(template_path, pdf_files):
    """
    Создает презентацию PowerPoint на основе шаблона, вставляя PDF-файлы на слайды.

    Args:
        template_path: Путь к файлу шаблона PowerPoint.
        pdf_files: Список путей к PDF-файлам.

    Returns:
        Путь к созданной презентации PowerPoint.
    """

    if not os.path.isdir('cache'):
        os.mkdir('cache')
    # Загружаем шаблон
    prs = Presentation(template_path)

    # Получаем слайды из шаблона
    slides = prs.slides

    # Проверяем, что в шаблоне достаточно слайдов для всех PDF-файлов
    # if len(slides) < len(pdf_files):
    #     raise ValueError("В шаблоне недостаточно слайдов для всех PDF-файлов.")

    # Вставляем PDF-файлы на слайды
    for i, pdf_file in enumerate(pdf_files):
        # Находим слайд с текстом-заменителем
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame and "{{" + pdf_file + "}}" == shape.text_frame.text:
                    images = convert_from_path(pdf_file)
                    # for i in images:
                    img_name = "cache/" + '.'.join(pdf_file.split('.')[:-1]) + '.png'
                    if len(images) > 1:
                        image = combine_images_horizontal(images)
                    else:
                        image = images[0]
                    image.save(img_name, format="png")

                    # Добавляем пустое место для изображения
                    with open(img_name, 'rb') as img_file:
                        # pic = slide.shapes.add_picture(
                        #     img_file,
                        #     left=Inches(prs.slide_width.inches * 0.1), top=Inches(prs.slide_height.inches * 0.2),
                        #     height=Inches(prs.slide_height.inches * 0.8)
                        # )
                        pic = slide.shapes.add_picture(
                            img_file,
                            left=Inches(0), top=Inches(0),
                            height=Inches(prs.slide_height.inches)
                        )
                        if pic.width > prs.slide_width:
                            pic.width = Inches(prs.slide_width.inches)
                        else:
                            pic.left = (prs.slide_width - pic.width) // 2


                    # Удаляем текст-заполнитель
                    shape.text_frame.text = ""
                    break

    # shutil.rmtree("cache")

    # Сохраняем презентацию
    output_path = "result.pptx"
    prs.save(output_path)

    return output_path


# Пример использования:
template_path = "./gen/template.pptx"
pdf_files = filter(lambda x: x.endswith('.pdf'), os.listdir())

output_path = create_presentation(template_path, pdf_files)

print(f"Презентация создана: {output_path}")
